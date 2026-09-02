"""Скрипт для генерации тестовых документов: PDF, DOCX, PPTX, XLSX."""

import os

# Папка для тестовых файлов
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "test_files")
os.makedirs(OUTPUT_DIR, exist_ok=True)


def create_xlsx():
    """Создаёт тестовый XLSX с несколькими листами."""
    from openpyxl import Workbook

    wb = Workbook()

    # Лист 1: FAQ по продукту
    ws1 = wb.active
    ws1.title = "FAQ Продукт"
    ws1.append(["Вопрос", "Ответ", "Категория"])
    ws1.append(["Как оформить возврат?", "Заполните форму на сайте в разделе 'Возвраты' в течение 14 дней", "Возвраты"])
    ws1.append(["Какие способы оплаты?", "Карта, СБП, наличные при получении", "Оплата"])
    ws1.append(["Сколько идёт доставка?", "По Москве — 1-2 дня, по России — 3-7 дней", "Доставка"])
    ws1.append(["Есть ли гарантия?", "Да, 12 месяцев на все товары", "Гарантия"])
    ws1.append(["Как отследить заказ?", "В личном кабинете в разделе 'Мои заказы'", "Доставка"])

    # Лист 2: Тарифы
    ws2 = wb.create_sheet("Тарифы")
    ws2.append(["Тариф", "Цена/мес", "Лимит запросов", "Поддержка"])
    ws2.append(["Базовый", "990 ₽", "1 000", "Email"])
    ws2.append(["Стандарт", "2 990 ₽", "10 000", "Email + чат"])
    ws2.append(["Про", "9 990 ₽", "100 000", "24/7 + менеджер"])
    ws2.append(["Корпоративный", "По запросу", "Безлимит", "Выделенная команда"])

    # Лист 3: Контакты
    ws3 = wb.create_sheet("Контакты")
    ws3.append(["Отдел", "Телефон", "Email", "Часы работы"])
    ws3.append(["Продажи", "+7 (495) 123-45-67", "sales@company.ru", "9:00-18:00"])
    ws3.append(["Поддержка", "+7 (495) 123-45-68", "support@company.ru", "Круглосуточно"])
    ws3.append(["Бухгалтерия", "+7 (495) 123-45-69", "billing@company.ru", "10:00-17:00"])

    path = os.path.join(OUTPUT_DIR, "test_data.xlsx")
    wb.save(path)
    print(f"[OK] XLSX: {path}")


def create_docx():
    """Создаёт тестовый DOCX с текстом и таблицами."""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading("Руководство пользователя — Сервис «Облако»", level=1)

    doc.add_paragraph(
        "Данное руководство описывает основные функции сервиса «Облако» "
        "и отвечает на часто задаваемые вопросы пользователей."
    )

    doc.add_heading("1. Системные требования", level=2)

    # Таблица 1: системные требования
    table1 = doc.add_table(rows=1, cols=3)
    table1.style = "Table Grid"
    hdr = table1.rows[0].cells
    hdr[0].text = "Параметр"
    hdr[1].text = "Минимальные"
    hdr[2].text = "Рекомендуемые"

    data1 = [
        ("ОС", "Windows 10 / macOS 12", "Windows 11 / macOS 14"),
        ("ОЗУ", "4 ГБ", "8 ГБ"),
        ("Диск", "500 МБ", "2 ГБ SSD"),
        ("Браузер", "Chrome 90+", "Chrome последней версии"),
        ("Интернет", "5 Мбит/с", "20 Мбит/с"),
    ]
    for param, min_val, rec_val in data1:
        row = table1.add_row().cells
        row[0].text = param
        row[1].text = min_val
        row[2].text = rec_val

    doc.add_paragraph("")
    doc.add_heading("2. Часто задаваемые вопросы", level=2)

    doc.add_paragraph(
        "Ниже представлены ответы на наиболее распространённые вопросы "
        "наших клиентов."
    )

    # Таблица 2: FAQ
    table2 = doc.add_table(rows=1, cols=2)
    table2.style = "Table Grid"
    hdr2 = table2.rows[0].cells
    hdr2[0].text = "Вопрос"
    hdr2[1].text = "Ответ"

    faq = [
        ("Как создать аккаунт?", "Нажмите 'Регистрация' на главной странице и заполните форму."),
        ("Забыл пароль, что делать?", "Нажмите 'Забыли пароль?' на странице входа. На вашу почту придёт ссылка для сброса."),
        ("Можно ли работать офлайн?", "Частично. Ранее загруженные файлы доступны без интернета, но синхронизация требует подключения."),
        ("Как поделиться файлом?", "Откройте файл → меню '...' → 'Поделиться' → скопируйте ссылку или введите email."),
    ]
    for q, a in faq:
        row = table2.add_row().cells
        row[0].text = q
        row[1].text = a

    doc.add_paragraph("")
    doc.add_paragraph(
        "Если вы не нашли ответ на свой вопрос, обратитесь в службу поддержки: support@cloud.ru"
    )

    path = os.path.join(OUTPUT_DIR, "test_manual.docx")
    doc.save(path)
    print(f"[OK] DOCX: {path}")


def create_pptx():
    """Создаёт тестовый PPTX с текстом и таблицами."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Слайд 1: Титульный
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Обзор продукта: Платформа «Аналитика»"
    slide1.placeholders[1].text = "Квартальный отчёт Q3 2026"

    # Слайд 2: Текст
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Ключевые метрики"
    tf = slide2.placeholders[1].text_frame
    tf.text = "За третий квартал 2026 года платформа показала стабильный рост:"
    tf.add_paragraph().text = "• Рост пользователей: +23%"
    tf.add_paragraph().text = "• Среднее время сессии: 12 минут"
    tf.add_paragraph().text = "• NPS: 72 балла"
    tf.add_paragraph().text = "• Uptime: 99.95%"

    # Слайд 3: Таблица
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    slide3.shapes.title.text = "Сравнение тарифов"

    rows, cols = 5, 4
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9.0)
    height = Inches(3.0)
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

    headers = ["Функция", "Free", "Pro", "Enterprise"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h

    data = [
        ("Хранилище", "5 ГБ", "100 ГБ", "Безлимит"),
        ("API-запросы/день", "100", "10 000", "Безлимит"),
        ("Пользователей", "1", "10", "Без ограничений"),
        ("Поддержка", "Форум", "Email + чат", "Выделенный менеджер"),
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # Слайд 4: Ещё таблица
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "Результаты по регионам"

    rows2, cols2 = 5, 3
    table2 = slide4.shapes.add_table(rows2, cols2, left, top, width, height).table
    for i, h in enumerate(["Регион", "Пользователей", "Выручка"]):
        table2.cell(0, i).text = h
    regions = [
        ("Москва", "45 000", "12.5 млн ₽"),
        ("Санкт-Петербург", "18 000", "4.8 млн ₽"),
        ("Новосибирск", "7 500", "1.9 млн ₽"),
        ("Другие", "29 500", "6.3 млн ₽"),
    ]
    for r, (reg, users, rev) in enumerate(regions, 1):
        table2.cell(r, 0).text = reg
        table2.cell(r, 1).text = users
        table2.cell(r, 2).text = rev

    path = os.path.join(OUTPUT_DIR, "test_presentation.pptx")
    prs.save(path)
    print(f"[OK] PPTX: {path}")


def create_pdf():
    """Создаёт тестовый PDF с текстом и таблицами через pdfplumber-совместимый формат."""
    # Используем простой подход через reportlab, если доступен,
    # иначе создаём PDF вручную через fpdf2
    try:
        from fpdf import FPDF
    except ImportError:
        print("[!] Dlya sozdaniya PDF nuzhen fpdf2: pip install fpdf2")
        print("   Propuskayu generaciyu PDF.")
        return

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Шрифт с поддержкой кириллицы — используем встроенный Helvetica
    # fpdf2 поддерживает Unicode
    pdf.add_page()

    # Заголовок
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "Spravochnik sotrudnika", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.ln(5)

    # Текст
    pdf.set_font("Helvetica", "", 11)
    pdf.multi_cell(0, 7,
        "Dannyj dokument soderzhit osnovnuyu informaciyu dlya novyh sotrudnikov "
        "kompanii. Pozhalujsta, oznakomtes' s tablicami nizhe."
    )
    pdf.ln(5)

    # Таблица 1: Отделы
    pdf.set_font("Helvetica", "B", 13)
    pdf.cell(0, 10, "1. Otdely kompanii", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)

    pdf.set_font("Helvetica", "B", 10)
    col_w = [50, 45, 50, 45]
    headers = ["Otdel", "Rukovoditel", "Email", "Etazh"]
    for i, h in enumerate(headers):
        pdf.cell(col_w[i], 8, h, border=1, align="C")
    pdf.ln()

    pdf.set_font("Helvetica", "", 10)
    rows = [
        ("Razrabotka", "Ivanov A.S.", "dev@co.ru", "3"),
        ("Marketing", "Petrova M.K.", "mkt@co.ru", "2"),
        ("Prodazhi", "Sidorov V.N.", "sales@co.ru", "1"),
        ("HR", "Kozlova E.A.", "hr@co.ru", "2"),
    ]
    for row in rows:
        for i, val in enumerate(row):
            pdf.cell(col_w[i], 7, val, border=1)
        pdf.ln()

    pdf.ln(8)

    # Текст между таблицами
    pdf.set_font("Helvetica", "", 11)
    pdf.multi_cell(0, 7,
        "Grafik raboty: ponedelnik-pyatnica, 9:00-18:00. "
        "Obedennyi pereryv: 13:00-14:00."
    )
    pdf.ln(5)

    # Таблица 2: Льготы
    pdf.set_font("Helvetica", "B", 13)
    pdf.cell(0, 10, "2. Lgoty sotrudnikov", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)

    pdf.set_font("Helvetica", "B", 10)
    col_w2 = [60, 65, 65]
    for i, h in enumerate(["Lgota", "Opisanie", "Usloviya"]):
        pdf.cell(col_w2[i], 8, h, border=1, align="C")
    pdf.ln()

    pdf.set_font("Helvetica", "", 9)
    benefits = [
        ("DMS", "Meditsinskaya strahovka", "Posle isp. sroka"),
        ("Fitnes", "Kompensatsiya 5000r/mes", "Vsem sotrudnikam"),
        ("Obuchenie", "Budget 50000r/god", "Soglasovanie s HR"),
        ("Udalyonka", "2 dnya v nedelyu", "Po soglasovaniyu"),
    ]
    for row in benefits:
        for i, val in enumerate(row):
            pdf.cell(col_w2[i], 7, val, border=1)
        pdf.ln()

    path = os.path.join(OUTPUT_DIR, "test_handbook.pdf")
    pdf.output(path)
    print(f"[OK] PDF: {path}")


if __name__ == "__main__":
    print(f"Generating test files in {OUTPUT_DIR}...\n")
    create_xlsx()
    create_docx()
    create_pptx()
    create_pdf()
    print(f"\n[DONE] Files in: {OUTPUT_DIR}")

