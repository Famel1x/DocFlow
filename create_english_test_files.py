"""Script to generate test files completely in English (XLSX, DOCX, PPTX, PDF scan).
Avoids Windows console Cyrillic encoding issues (diamond/square question marks)
and provides clean test suites for all 6 rule modes and OCR.
"""

import os
import io
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill, Border, Side
import docx
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
import fitz

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "test_files")
os.makedirs(OUTPUT_DIR, exist_ok=True)


# ══════════════════════════════════════════════════════════════════════
# 1. English XLSX with Merged Cells & Multi-level Headers
# ══════════════════════════════════════════════════════════════════════
def create_english_xlsx():
    wb = Workbook()

    # Sheet 1: Grouped Pricing Catalog (Merged headers & rows)
    ws1 = wb.active
    ws1.title = "Product_Catalog"

    # Multi-level Header
    ws1.append(["Product Category", "", "Standard Tier", "", "Enterprise Tier", ""])
    ws1.merge_cells("A1:B1")
    ws1.merge_cells("C1:D1")
    ws1.merge_cells("E1:F1")

    ws1.append(["Category", "Item Name", "Price USD", "Warranty", "Price USD", "Warranty"])

    # Data Rows
    data1 = [
        ["Hardware", "Cloud Server X1", "1,200", "12 mo", "1,800", "36 mo"],
        ["Hardware", "Edge Router R40", "450", "12 mo", "750", "36 mo"],
        ["Hardware", "Managed Switch S24", "650", "12 mo", "950", "36 mo"],
        ["Software", "Database License", "300", "Annual", "600", "Perpetual"],
        ["Software", "AI Assistant API", "500", "Monthly", "1,200", "Annual"],
        ["Software", "Security Suite Pro", "250", "Monthly", "500", "Annual"],
    ]
    for row in data1:
        ws1.append(row)

    ws1.merge_cells("A3:A5")
    ws1.merge_cells("A6:A8")

    # Sheet 2: Access Control Matrix (Role x Module)
    ws2 = wb.create_sheet("Access_Matrix")
    ws2.append(["User Role", "CRM Module", "Billing Portal", "Warehouse", "Analytics Dashboard", "System Settings"])
    ws2.append(["Intern", "Read Only", "No Access", "Read Only", "No Access", "No Access"])
    ws2.append(["Sales Manager", "Full Access", "Read Only", "Create Orders", "Sales Reports", "No Access"])
    ws2.append(["Support Lead", "Full Access", "No Access", "Read Only", "Ticket Stats", "Basic Settings"])
    ws2.append(["Finance Officer", "Read Only", "Full Access", "Audit Logs", "Financial Reports", "No Access"])
    ws2.append(["Super Admin", "Full Access", "Full Access", "Full Access", "Full Access", "Full Access"])

    # Sheet 3: Standard Customer Support FAQ (2 columns)
    ws3 = wb.create_sheet("Support_FAQ")
    ws3.append(["Question", "Answer", "Category"])
    faq_data = [
        ["How do I reset my account password?", "Click 'Forgot Password' on the login screen and follow the email link.", "Account"],
        ["What payment methods are supported?", "We accept Visa, MasterCard, Wire Transfer, and PayPal.", "Billing"],
        ["Can I upgrade my subscription plan anytime?", "Yes, navigate to Settings -> Billing -> Change Plan.", "Billing"],
        ["How to configure custom webhook alerts?", "Go to Developer Portal -> Webhooks -> Add Endpoint.", "API & Dev"],
        ["Where can I find my API secret keys?", "Your API keys are located under Settings -> Security -> API Tokens.", "Security"]
    ]
    for r in faq_data:
        ws3.append(r)

    path = os.path.join(OUTPUT_DIR, "english_demo.xlsx")
    wb.save(path)
    print(f"[OK] English XLSX created: {path}")


# ══════════════════════════════════════════════════════════════════════
# 2. English DOCX with System Specs & FAQ Tables
# ══════════════════════════════════════════════════════════════════════
def set_cell_merge_grid_span(cell, span):
    tcPr = cell._tc.get_or_add_tcPr()
    gridSpan = OxmlElement('w:gridSpan')
    gridSpan.set(qn('w:val'), str(span))
    tcPr.append(gridSpan)


def set_cell_v_merge(cell, val):
    tcPr = cell._tc.get_or_add_tcPr()
    vMerge = OxmlElement('w:vMerge')
    if val:
        vMerge.set(qn('w:val'), val)
    tcPr.append(vMerge)


def create_english_docx():
    doc = docx.Document()
    doc.add_heading("Cloud Infrastructure User Manual", level=1)
    doc.add_paragraph("This technical guide covers specifications and frequently asked questions.")

    # Table 1: System Requirements (3 columns)
    doc.add_heading("1. Minimum and Recommended System Requirements", level=2)
    t1 = doc.add_table(rows=1, cols=3)
    t1.style = "Table Grid"
    hdr1 = t1.rows[0].cells
    hdr1[0].text = "Component"
    hdr1[1].text = "Minimum Spec"
    hdr1[2].text = "Recommended Spec"

    specs = [
        ("Operating System", "Ubuntu 22.04 LTS / Windows Server 2022", "Ubuntu 24.04 LTS"),
        ("CPU Cores", "4 vCPU @ 2.4 GHz", "16 vCPU @ 3.2 GHz"),
        ("System Memory", "16 GB DDR4 ECC", "64 GB DDR5 ECC"),
        ("Disk Storage", "100 GB NVMe SSD", "1 TB NVMe RAID-10"),
        ("Network Bandwidth", "1 Gbps Uplink", "10 Gbps Redundant Uplink"),
    ]
    for comp, min_s, rec_s in specs:
        r = t1.add_row().cells
        r[0].text = comp
        r[1].text = min_s
        r[2].text = rec_s

    # Table 2: FAQ Table (Question, Answer, Severity)
    doc.add_heading("2. Troubleshooting FAQ", level=2)
    t2 = doc.add_table(rows=1, cols=3)
    t2.style = "Table Grid"
    hdr2 = t2.rows[0].cells
    hdr2[0].text = "Question"
    hdr2[1].text = "Resolution"
    hdr2[2].text = "Priority"

    faqs = [
        ("Server fails to restart after kernel update", "Boot into recovery console and select the prior kernel image.", "High"),
        ("Connection timeout when accessing HTTPS port 443", "Verify firewall security group rules allow inbound traffic on 443.", "Critical"),
        ("Disk usage alert triggered above 85 percent", "Run log rotation utility and clear cached package archives.", "Medium"),
        ("SSL certificate expiring within 7 days", "Run certbot renew or upload renewed PEM certificate chain.", "High"),
    ]
    for q, a, p in faqs:
        r = t2.add_row().cells
        r[0].text = q
        r[1].text = a
        r[2].text = p

    path = os.path.join(OUTPUT_DIR, "english_demo.docx")
    doc.save(path)
    print(f"[OK] English DOCX created: {path}")


# ══════════════════════════════════════════════════════════════════════
# 3. English PPTX Presentation with Slides & Tables
# ══════════════════════════════════════════════════════════════════════
def create_english_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Slide 1: SLA Service Levels Table
    s1 = prs.slides.add_slide(blank_layout)
    txBox = s1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(0.8))
    txBox.text_frame.text = "Service Level Agreement (SLA) Matrix"

    t_shape = s1.shapes.add_table(5, 4, Inches(1), Inches(1.5), Inches(11.3), Inches(4.5))
    table = t_shape.table

    headers = ["Service Tier", "Uptime Guarantee", "Response SLA", "Dedicated Support"]
    for c_idx, h in enumerate(headers):
        table.cell(0, c_idx).text = h

    rows = [
        ["Bronze Support", "99.0% Uptime", "Within 24 Business Hours", "Email Only"],
        ["Silver Support", "99.5% Uptime", "Within 8 Business Hours", "Email + Chat"],
        ["Gold Support", "99.9% Uptime", "Within 2 Business Hours", "24/7 Phone + Slack"],
        ["Mission Critical", "99.99% Uptime", "15 Minutes Guaranteed", "Dedicated TAM + War Room"]
    ]
    for r_idx, row in enumerate(rows, 1):
        for c_idx, val in enumerate(row):
            table.cell(r_idx, c_idx).text = val

    path = os.path.join(OUTPUT_DIR, "english_demo.pptx")
    prs.save(path)
    print(f"[OK] English PPTX created: {path}")


# ══════════════════════════════════════════════════════════════════════
# 4. English Scanned PDF Document (For Testing OCR without artifacts)
# ══════════════════════════════════════════════════════════════════════
def create_english_scanned_pdf():
    img = Image.new('RGB', (900, 650), color=(255, 255, 255))
    d = ImageDraw.Draw(img)

    # Document Title
    d.text((50, 30), "SCANNED HARDWARE ASSET REPORT (CONFIDENTIAL)", fill=(20, 20, 20))
    d.text((50, 55), "Warehouse Inventory - Audited Q3", fill=(100, 100, 100))

    headers = ["Asset ID", "Device Description", "Serial Number", "Unit Cost", "Status"]
    rows = [
        ["SRV-101", "Dual Xeon Compute Node", "SN-88239-X", "$3,450", "In Service"],
        ["SW-204", "Arista 32-Port 100G Switch", "SN-44102-A", "$5,200", "In Service"],
        ["STR-502", "PureStorage NVMe Array 50TB", "SN-99120-P", "$18,500", "Operational"],
        ["FW-301", "Palo Alto Next-Gen Firewall", "SN-12093-F", "$8,900", "Under Maintenance"],
        ["PDU-04", "APC Monitored Rack PDU", "SN-77231-E", "$850", "Spare"]
    ]

    col_widths = [100, 260, 150, 120, 150]
    y = 110

    # Draw Header Row
    x = 50
    for idx, h in enumerate(headers):
        w = col_widths[idx]
        d.rectangle([x, y, x + w, y + 40], fill=(235, 240, 250), outline=(80, 100, 140), width=2)
        d.text((x + 10, y + 12), h, fill=(10, 20, 50))
        x += w
    y += 40

    # Draw Data Rows
    for r in rows:
        x = 50
        for idx, val in enumerate(r):
            w = col_widths[idx]
            d.rectangle([x, y, x + w, y + 38], outline=(160, 170, 185), width=1)
            d.text((x + 10, y + 11), val, fill=(30, 30, 30))
            x += w
        y += 38

    # Save PNG and build non-searchable scanned PDF
    png_path = os.path.join(OUTPUT_DIR, "scanned_english_asset.png")
    img.save(png_path)

    pdf_path = os.path.join(OUTPUT_DIR, "scanned_english_document.pdf")
    doc = fitz.open()
    page = doc.new_page(width=900, height=650)
    page.insert_image(page.rect, filename=png_path)
    doc.save(pdf_path)
    print(f"[OK] English Scanned PDF created: {pdf_path}")


if __name__ == "__main__":
    create_english_xlsx()
    create_english_docx()
    create_english_pptx()
    create_english_scanned_pdf()
    print("\n[SUCCESS] All English test files successfully generated in test_files/ folder!")
