"""Парсер PPTX с поддержкой объединённых ячеек."""

from pptx import Presentation
from typing import Any
from app.parsers.base import BaseParser


class PPTXParser(BaseParser):
    def parse(self) -> list[dict[str, Any]]:
        content = []
        prs = Presentation(self.file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            content.append({"type": "text", "value": f"--- Слайд {slide_num} ---"})

            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    rows = len(tbl.rows)
                    cols = len(tbl.columns)

                    # Строим плоскую сетку
                    grid = [[""] * cols for _ in range(rows)]
                    merges = []

                    for r_idx, row in enumerate(tbl.rows):
                        for c_idx in range(cols):
                            cell = tbl.cell(r_idx, c_idx)
                            grid[r_idx][c_idx] = cell.text_frame.text.strip()

                            # Проверяем объединение через span
                            if cell.is_merge_origin:
                                span_h = cell.span_width
                                span_v = cell.span_height
                                text = grid[r_idx][c_idx]

                                # Заполняем все ячейки в объединении
                                for dr in range(span_v):
                                    for dc in range(span_h):
                                        if dr == 0 and dc == 0:
                                            continue
                                        if r_idx + dr < rows and c_idx + dc < cols:
                                            grid[r_idx + dr][c_idx + dc] = text

                                if span_h > 1 or span_v > 1:
                                    merges.append({
                                        "min_row": r_idx, "min_col": c_idx,
                                        "max_row": r_idx + span_v - 1,
                                        "max_col": c_idx + span_h - 1,
                                    })

                    item = {"type": "table", "value": grid}
                    if merges:
                        item["merges"] = merges
                    content.append(item)

                elif shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        content.append({"type": "text", "value": text})

        return content
