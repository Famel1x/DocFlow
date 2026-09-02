"""Скрипт генерации сложных тестовых файлов с объединенными ячейками (XLSX, DOCX, PPTX)."""

import os
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill, Border, Side
import docx
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "test_files")
os.makedirs(OUTPUT_DIR, exist_ok=True)


def set_cell_background_and_border(cell):
    thin = Side(border_style="thin", color="CCCCCC")
    cell.border = Border(top=thin, left=thin, right=thin, bottom=thin)


def create_complex_xlsx():
    wb = Workbook()
    
    # Лист 1: Сложный многоуровневый прайс-лист с объединением по вертикали и горизонтали
    ws1 = wb.active
    ws1.title = "Сложный_Прайс"
    
    # Заголовки с двухуровневой шапкой
    # Строка 1: Категория (span 2 cols), Базовые условия (span 2 cols), Премиум (span 2 cols)
    ws1.append(["Категория товара", "", "Базовый пакет", "", "Премиум пакет", ""])
    ws1.merge_cells("A1:B1")
    ws1.merge_cells("C1:D1")
    ws1.merge_cells("E1:F1")
    
    # Строка 2: подзаголовки
    ws1.append(["Группа", "Наименование", "Цена (руб)", "Срок гарантии", "Цена (руб)", "Срок гарантии"])
    
    # Данные с вертикальным объединением по Категории (Группа)
    data = [
        # Бытовая техника (строки 3-5, A3:A5 объединено)
        ["Бытовая техника", "Стиральная машина Alfa", "35 000", "12 мес", "45 000", "36 мес"],
        ["Бытовая техника", "Холодильник FrostPro", "52 000", "12 мес", "68 000", "36 мес"],
        ["Бытовая техника", "Посудомоечная машина CleanX", "41 000", "12 мес", "55 000", "36 мес"],
        # Компьютеры (строки 6-8, A6:A8 объединено)
        ["Компьютеры", "Ноутбук WorkBook 14", "65 000", "12 мес", "85 000", "24 мес"],
        ["Компьютеры", "Системный блок GamingPro", "110 000", "24 мес", "145 000", "36 мес"],
        ["Компьютеры", "Монитор UltraView 27", "28 000", "24 мес", "38 000", "36 мес"],
    ]
    
    for row in data:
        ws1.append(row)
        
    ws1.merge_cells("A3:A5")
    ws1.merge_cells("A6:A8")
    
    # Лист 2: Сводная матрица (перекрестная таблица)
    ws2 = wb.create_sheet("Матрица_Доступов")
    ws2.append(["Роль сотрудника", "Модуль CRM", "Бухгалтерия", "Склад", "Аналитика", "Настройки системы"])
    ws2.append(["Стажер", "Чтение", "Нет доступа", "Чтение", "Нет доступа", "Нет доступа"])
    ws2.append(["Менеджер продаж", "Полный доступ", "Нет доступа", "Создание заказов", "Отчеты по продажам", "Нет доступа"])
    ws2.append(["Бухгалтер", "Чтение клиентов", "Полный доступ", "Чтение остатков", "Финансовые отчеты", "Нет доступа"])
    ws2.append(["Кладовщик", "Нет доступа", "Нет доступа", "Полный доступ", "Остатки", "Нет доступа"])
    ws2.append(["Администратор", "Полный доступ", "Полный доступ", "Полный доступ", "Полный доступ", "Полный доступ"])
    
    path = os.path.join(OUTPUT_DIR, "complex_merged.xlsx")
    wb.save(path)
    print(f"[OK] Complex XLSX: {path}")


def create_complex_docx():
    doc = docx.Document()
    doc.add_heading("Спецификация услуг и тарифов (Сложная таблица)", level=1)
    
    doc.add_paragraph("Ниже представлена таблица тарификации с вертикальными и горизонтальными объединениями ячеек:")
    
    # Таблица: 6 строк, 4 колонки
    table = doc.add_table(rows=6, cols=4)
    table.style = 'Table Grid'
    
    # Заголовок (строка 0): Объединяем колонки 2 и 3 в "Параметры обслуживания"
    cell_hdr_cat = table.cell(0, 0)
    cell_hdr_cat.text = "Категория"
    
    cell_hdr_srv = table.cell(0, 1)
    cell_hdr_srv.text = "Услуга"
    
    # Horizontal merge cell(0,2) and cell(0,3)
    cell_params = table.cell(0, 2)
    cell_params.merge(table.cell(0, 3))
    cell_params.text = "Условия и Лимиты"
    
    # Строка 1: Подзаголовки
    table.cell(1, 0).text = "Категория"
    table.cell(1, 1).text = "Услуга"
    table.cell(1, 2).text = "SLA (время реакции)"
    table.cell(1, 3).text = "Стоимость / мес"
    
    # Строки 2 и 3: Категория "Инфраструктура" объединена по вертикали
    table.cell(2, 0).text = "Инфраструктура"
    table.cell(2, 1).text = "Аренда выделенного сервера"
    table.cell(2, 2).text = "До 15 минут"
    table.cell(2, 3).text = "15 000 руб."
    
    table.cell(3, 0).text = "Инфраструктура"
    table.cell(3, 1).text = "Облачное хранилище S3"
    table.cell(3, 2).text = "До 1 часа"
    table.cell(3, 3).text = "3 000 руб."
    
    # Вертикальное объединение строки 2 и 3 в первой колонке
    table.cell(2, 0).merge(table.cell(3, 0))
    
    # Строки 4 и 5: Категория "Безопасность" объединена по вертикали
    table.cell(4, 0).text = "Безопасность"
    table.cell(4, 1).text = "DDoS защита L3-L7"
    table.cell(4, 2).text = "Мгновенно (авто)"
    table.cell(4, 3).text = "8 000 руб."
    
    table.cell(5, 0).text = "Безопасность"
    table.cell(5, 1).text = "Аудит уязвимостей"
    table.cell(5, 2).text = "1 рабочий день"
    table.cell(5, 3).text = "25 000 руб."
    
    table.cell(4, 0).merge(table.cell(5, 0))
    
    path = os.path.join(OUTPUT_DIR, "complex_merged.docx")
    doc.save(path)
    print(f"[OK] Complex DOCX: {path}")


def create_complex_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Сводный план этапов проекта"
    
    rows, cols = 5, 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(3.5)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table
    
    # Заголовок (строка 0): 0: Фаза, 1: Задача, 2-3: Сроки и Ответственный (объединено по горизонтали)
    tbl.cell(0, 0).text = "Фаза проекта"
    tbl.cell(0, 1).text = "Основная задача"
    
    # Merge cells (0,2) and (0,3)
    tbl.cell(0, 2).merge(tbl.cell(0, 3))
    tbl.cell(0, 2).text = "Исполнение и контроль"
    
    # Строка 1: Фаза 1 (вертикально объединена со строкой 2)
    tbl.cell(1, 0).text = "Этап 1: Проектирование"
    tbl.cell(1, 1).text = "Сбор бизнес-требований"
    tbl.cell(1, 2).text = "Срок: 2 недели"
    tbl.cell(1, 3).text = "Аналитик: Иванов"
    
    tbl.cell(2, 0).text = "Этап 1: Проектирование"
    tbl.cell(2, 1).text = "Архитектурный проект"
    tbl.cell(2, 2).text = "Срок: 3 недели"
    tbl.cell(2, 3).text = "Архитектор: Смирнов"
    
    # Merge vertical (1,0) and (2,0)
    tbl.cell(1, 0).merge(tbl.cell(2, 0))
    
    # Строка 3: Фаза 2 (вертикально со строкой 4)
    tbl.cell(3, 0).text = "Этап 2: Разработка"
    tbl.cell(3, 1).text = "Создание backend API"
    tbl.cell(3, 2).text = "Срок: 4 недели"
    tbl.cell(3, 3).text = "Разработчик: Кузнецов"
    
    tbl.cell(4, 0).text = "Этап 2: Разработка"
    tbl.cell(4, 1).text = "Верстка и фронтенд интерфейса"
    tbl.cell(4, 2).text = "Срок: 3 недели"
    tbl.cell(4, 3).text = "Фронтенд: Попов"
    
    tbl.cell(3, 0).merge(tbl.cell(4, 0))
    
    path = os.path.join(OUTPUT_DIR, "complex_merged.pptx")
    prs.save(path)
    print(f"[OK] Complex PPTX: {path}")


if __name__ == "__main__":
    print(f"Generating complex test files with merged cells in {OUTPUT_DIR}...\n")
    create_complex_xlsx()
    create_complex_docx()
    create_complex_pptx()
    print(f"\n[DONE] Complex test files created successfully!")
